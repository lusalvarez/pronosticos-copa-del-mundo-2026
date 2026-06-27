from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Charger le template
template_path = r"C:\Users\059758706\Desktop\pronostics-coupe-du-monde\template power point\Do not use Horizontal template givaudan oral.pptx"
prs = Presentation(template_path)

print(f"Template charge avec succes!")
print(f"Nombre de layouts disponibles: {len(prs.slide_layouts)}")

# Afficher les layouts disponibles
for i, layout in enumerate(prs.slide_layouts):
    print(f"Layout {i}: {layout.name}")

# Données des matchs extraites de l'image
matches = [
    {
        "id": 24,
        "homeTeam": "Chequia",
        "awayTeam": "Sudafrica",
        "date": "2026-06-18T18:00",
        "stage": "Fase de grupos - Grupo A",
        "prediction": {"away": 1, "home": 1, "firstGoal": ""}
    },
    {
        "id": 25,
        "homeTeam": "Suiza",
        "awayTeam": "Bosnia y Herzegovina",
        "date": "2026-06-18T21:00",
        "stage": "Fase de grupos - Grupo B",
        "prediction": {"away": 1, "home": 1, "firstGoal": ""}
    },
    {
        "id": 26,
        "homeTeam": "Canada",
        "awayTeam": "Catar",
        "date": "2026-06-19T00:00",
        "stage": "Fase de grupos - Grupo B",
        "prediction": {"home": 1, "away": "", "firstGoal": ""}
    }
]

# Slide 1: Page de titre (utiliser le premier layout disponible)
try:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Essayer de trouver et remplir le titre
    title_found = False
    subtitle_found = False
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        # Essayer d'accéder au text_frame
        try:
            text_frame = shape.text_frame
            
            # Si c'est un placeholder
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1:  # Title
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = "Pronostics Coupe du Monde 2026"
                    p.font.size = Pt(40)
                    p.font.bold = True
                    title_found = True
                    print("Titre principal rempli")
                elif phf.type == 2:  # Subtitle
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = "Phase de Groupes - Matchs a venir"
                    p.font.size = Pt(24)
                    subtitle_found = True
                    print("Sous-titre rempli")
            # Si ce n'est pas un placeholder mais qu'on n'a pas encore trouvé de titre
            elif not title_found and text_frame.text == "":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = "Pronostics Coupe du Monde 2026"
                p.font.size = Pt(40)
                p.font.bold = True
                title_found = True
                print("Titre rempli dans zone de texte")
        except Exception as e:
            print(f"Erreur avec shape: {e}")
            continue
    
    if not title_found:
        print("ATTENTION: Titre non trouve, ajout manuel")
        # Ajouter manuellement une zone de texte pour le titre
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Pronostics Coupe du Monde 2026"
        p.font.size = Pt(40)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    print("Slide de titre creee")
except Exception as e:
    print(f"Erreur lors de la creation de la slide de titre: {e}")

# Créer une slide pour chaque match
for i, match in enumerate(matches, 1):
    try:
        # Utiliser un layout avec titre et contenu (généralement layout 1 ou 2)
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Remplir le titre
        title_filled = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == 1 and not title_filled:  # Title
                    shape.text = f"Match {match['id']}: {match['homeTeam']} vs {match['awayTeam']}"
                    title_filled = True
                    
                    # Formater le titre
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = Pt(28)
                        paragraph.font.bold = True
        
        # Ajouter le contenu dans une zone de texte
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Équipe domicile
        p = tf.paragraphs[0]
        p.text = f"Equipe a domicile: {match['homeTeam']}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(12)
        
        # Équipe extérieure
        p = tf.add_paragraph()
        p.text = f"Equipe a l'exterieur: {match['awayTeam']}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(16)
        
        # Date
        p = tf.add_paragraph()
        date_str = match['date'].replace('T', ' a ')
        p.text = f"Date: {date_str}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        
        # Phase
        p = tf.add_paragraph()
        p.text = f"Phase: {match['stage']}"
        p.font.size = Pt(18)
        p.space_after = Pt(16)
        
        # Pronostic
        pred = match['prediction']
        home_score = pred.get('home', '?')
        away_score = pred.get('away', '?')
        
        p = tf.add_paragraph()
        p.text = f"Pronostic: {home_score} - {away_score}"
        p.font.size = Pt(22)
        p.font.bold = True
        
        print(f"Slide {i} creee pour le match {match['id']}")
        
    except Exception as e:
        print(f"Erreur lors de la creation de la slide {i}: {e}")

# Slide finale: Résumé
try:
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.type == 1:  # Title
                shape.text = "Resume des Pronostics"
                break
    
    # Contenu
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Pronostics pour les matchs a venir:\n"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(12)
    
    for match in matches:
        pred = match['prediction']
        home_score = pred.get('home', '?')
        away_score = pred.get('away', '?')
        
        p = tf.add_paragraph()
        p.text = f"{match['homeTeam']} {home_score} - {away_score} {match['awayTeam']}"
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    print("Slide de resume creee")
except Exception as e:
    print(f"Erreur lors de la creation de la slide de resume: {e}")

# Sauvegarder la présentation
import datetime
timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
output_path = f"Pronostics_Coupe_du_Monde_2026_{timestamp}.pptx"
prs.save(output_path)
print(f"\nPresentation creee avec succes: {output_path}")
print(f"Nombre total de slides: {len(prs.slides)}")

# Made with Bob
