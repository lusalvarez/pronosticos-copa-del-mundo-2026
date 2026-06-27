from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Données des matchs extraites de l'image
matches = [
    {
        "id": 24,
        "homeTeam": "Chequia",
        "awayTeam": "Sudáfrica",
        "date": "2026-06-18T18:00",
        "stage": "Fase de grupos - Grupo A",
        "prediction": {"away": 1, "home": 1, "firstGoal": ""}
    },
    {
        "id": 25,
        "homeTeam": "Suiza",
        "awayTeam": "Bosnia y Herzegovina",
        "date": "2026-06-18T21:00",
        "stage": "Fase de grupos - Grupo B",
        "prediction": {"away": 1, "home": 1, "firstGoal": ""}
    },
    {
        "id": 26,
        "homeTeam": "Canadá",
        "awayTeam": "Catar",
        "date": "2026-06-19T00:00",
        "stage": "Fase de grupos - Grupo B",
        "prediction": {"home": 1, "away": "", "firstGoal": ""}
    }
]

# Slide 1: Page de titre
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Pronostics Coupe du Monde 2026"
subtitle.text = "Phase de Groupes - Matchs à venir"

# Formater le titre
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# Slide 2: Vue d'ensemble
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Vue d'ensemble des matchs"

# Ajouter un cadre de texte
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = f"Nombre total de matchs: {len(matches)}\n\n"
p.font.size = Pt(18)

for match in matches:
    p = tf.add_paragraph()
    p.text = f"• Match {match['id']}: {match['homeTeam']} vs {match['awayTeam']}"
    p.font.size = Pt(16)
    p.level = 0

# Créer une slide pour chaque match
for i, match in enumerate(matches, 1):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = f"Match {match['id']}: {match['homeTeam']} vs {match['awayTeam']}"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Informations du match
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(7)
    height = Inches(4.5)
    info_box = slide.shapes.add_textbox(left, top, width, height)
    info_frame = info_box.text_frame
    info_frame.word_wrap = True
    
    # Équipe domicile
    p = info_frame.paragraphs[0]
    p.text = f"🏠 Équipe à domicile: {match['homeTeam']}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(12)
    
    # Équipe extérieure
    p = info_frame.add_paragraph()
    p.text = f"✈️  Équipe à l'extérieur: {match['awayTeam']}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(20)
    
    # Date
    p = info_frame.add_paragraph()
    date_str = match['date'].replace('T', ' à ')
    p.text = f"📅 Date: {date_str}"
    p.font.size = Pt(20)
    p.space_after = Pt(12)
    
    # Phase
    p = info_frame.add_paragraph()
    p.text = f"🏆 Phase: {match['stage']}"
    p.font.size = Pt(20)
    p.space_after = Pt(20)
    
    # Pronostic
    pred = match['prediction']
    home_score = pred.get('home', '?')
    away_score = pred.get('away', '?')
    
    p = info_frame.add_paragraph()
    p.text = f"🎯 Pronostic: {home_score} - {away_score}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    
    if pred.get('firstGoal'):
        p = info_frame.add_paragraph()
        p.text = f"⚽ Premier but: {pred['firstGoal']}"
        p.font.size = Pt(18)

# Slide finale: Résumé
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Résumé des Pronostics"

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Pronostics pour les matchs à venir:\n\n"
p.font.size = Pt(20)
p.font.bold = True

for match in matches:
    pred = match['prediction']
    home_score = pred.get('home', '?')
    away_score = pred.get('away', '?')
    
    p = tf.add_paragraph()
    p.text = f"• {match['homeTeam']} {home_score} - {away_score} {match['awayTeam']}"
    p.font.size = Pt(18)
    p.space_after = Pt(8)

# Sauvegarder la présentation
prs.save('Pronostics_Coupe_du_Monde_2026.pptx')
print("Presentation creee avec succes: Pronostics_Coupe_du_Monde_2026.pptx")

# Made with Bob
